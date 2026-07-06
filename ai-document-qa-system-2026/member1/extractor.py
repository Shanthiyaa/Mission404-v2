"""
member1/extractor.py — Document extraction and chunking pipeline
Uses PyMuPDF for PDF (fast, reliable).
Docling removed — was causing install/runtime failures.
"""

import os
import re
import json
import hashlib
import logging
from pathlib import Path
from datetime import datetime
import zipfile
import tempfile

# ── LangChain splitters ───────────────────────────────────────────────────────
from langchain_text_splitters import (
    RecursiveCharacterTextSplitter,
    MarkdownHeaderTextSplitter,
)
from langchain_core.documents import Document

# ── Language detection (optional) ────────────────────────────────────────────
try:
    from langdetect import detect as _detect_lang
    _LANGDETECT_OK = True
except ImportError:
    _LANGDETECT_OK = False

# ── Shared config ─────────────────────────────────────────────────────────────
import sys
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from config import (
    SUPPORTED_EXTENSIONS, DOC_TYPE_MAP, CHUNK_CONFIG, CHUNKS_JSON_PATH
)

log = logging.getLogger(__name__)

_MD_HEADERS = [("#", "h1"), ("##", "h2"), ("###", "h3")]

# FIX BUG 3: pattern that identifies Table-of-Contents lines (dots + page ref)
_TOC_LINE_RE = re.compile(r'\.{4,}\s*[\dxliv]', re.I)


def _is_toc_chunk(text: str) -> bool:
    """
    Return True if this chunk is dominated by TOC lines.
    TOC chunks waste retrieval slots — they look relevant (they mention topic names)
    but contain zero instructional content.
    A chunk is considered TOC when ≥ 40% of its non-empty lines are TOC entries.
    """
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    if not lines:
        return False
    toc_count = sum(1 for l in lines if _TOC_LINE_RE.search(l))
    return (toc_count / len(lines)) >= 0.40


_NOISE_RE = [
    (re.compile(r"Page\s+\d+\s+of\s+\d+", re.I), ""),
    (re.compile(r"©\s*\d{4}[^\n]*",        re.I), ""),
    (re.compile(r"Confidential[^\n]*",      re.I), ""),
    (re.compile(r"Proprietary[^\n]*",       re.I), ""),
    (re.compile(r"\ufb01"),                        "fi"),
    (re.compile(r"\ufb02"),                        "fl"),
    (re.compile(r"[\u2018\u2019]"),                "'"),
    (re.compile(r"[\u201c\u201d]"),                '"'),
    (re.compile(r"[\u2013\u2014]"),                "-"),
    (re.compile(r"\u00a0"),                        " "),
    (re.compile(r"\x00"),                          ""),
    (re.compile(r"\n{3,}"),                        "\n\n"),
    (re.compile(r"[ \t]{2,}"),                     " "),
]


# ══════════════════════════════════════════════════════════════════════════════
#  UTILITIES
# ══════════════════════════════════════════════════════════════════════════════

def detect_doc_type(filename: str) -> str:
    name = filename.lower()
    for dtype, keywords in DOC_TYPE_MAP.items():
        if any(k in name for k in keywords):
            return dtype
    return "unknown"


def detect_language(text: str) -> str:
    if not _LANGDETECT_OK or len(text.strip()) < 40:
        return "unknown"
    try:
        return _detect_lang(text)
    except Exception:
        return "unknown"


def chunk_id(text: str, source: str, idx: int) -> str:
    return hashlib.md5(f"{source}::{idx}::{text[:100]}".encode()).hexdigest()[:14]


def clean_text(text: str) -> str:
    for pattern, replacement in _NOISE_RE:
        text = pattern.sub(replacement, text)
    return text.strip()


def calculate_file_hash(filepath: str) -> str:
    hasher = hashlib.sha256()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            hasher.update(chunk)
    return hasher.hexdigest()


def is_duplicate(filepath: str) -> tuple[bool, str]:
    """
    Check if the file content already exists in the knowledge base.
    Compares the SHA-256 hash of filepath with all files in UPLOAD_DIR.
    """
    from config import UPLOAD_DIR
    target_path = Path(filepath).resolve()
    if not target_path.exists():
        return False, ""
    
    target_hash = calculate_file_hash(str(target_path))
    upload_dir_path = Path(UPLOAD_DIR).resolve()
    
    if not upload_dir_path.exists():
        return False, ""
        
    for p in upload_dir_path.iterdir():
        if p.is_file() and p.resolve() != target_path:
            try:
                if calculate_file_hash(str(p)) == target_hash:
                    return True, p.name
            except Exception as e:
                log.warning(f"Failed to hash {p.name}: {e}")
                
    return False, ""


# ══════════════════════════════════════════════════════════════════════════════
#  EXTRACTION — PyMuPDF for PDF (page-aware, fast, no ML models needed)
# ══════════════════════════════════════════════════════════════════════════════

def extract_pdf(filepath: str) -> dict:
    """
    Extract PDF using PyMuPDF (fitz).
    Returns page-aware segments with real page numbers.
    """
    import fitz  # PyMuPDF

    p        = Path(filepath)
    filename = p.name
    doc_type = detect_doc_type(filename)

    log.info(f"Extracting PDF: {filename}")

    page_segments = []
    tables        = []
    current_section = ""

    doc = fitz.open(str(p))
    total_pages = len(doc)
    log.info(f"  {filename}: {total_pages} pages")

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text")
        if not text.strip():
            continue

        text = clean_text(text)
        lines = text.split("\n")

        for line in lines:
            line = line.strip()
            # FIX BUG 2: Better heading detection for ALE technical manuals.
            # Old logic caught only ALL-CAPS or lines ending in ":" — missing
            # most real headings like "Configuring VLANs" and catching noise like
            # "For example". New rules:
            #  1. All-caps, 10–80 chars, no digits-only content → chapter heading
            #  2. Title-case line, 10–80 chars, no trailing punctuation → section heading
            #  3. Skip lines that look like TOC entries (........number pattern)
            if not line or len(line) < 8 or len(line) > 80:
                continue
            if re.search(r'\.{3,}\s*\d', line):          # TOC line — skip
                continue
            if line.isupper() and re.search(r'[A-Z]', line):
                current_section = line                    # ALL-CAPS chapter heading
            elif (
                line[0].isupper()                         # starts with capital
                and not line.endswith((",", ";", ".", ":", "-"))  # no trailing punct
                and not re.match(r'^[\d\s\-\.]+$', line)  # not a number-only line
                and sum(1 for w in line.split() if w and w[0].isupper()) >= 2
            ):
                current_section = line                    # Title Case section heading

        page_segments.append({
            "text":    text,
            "page":    page_num,
            "section": current_section,
        })

    doc.close()

    full_text = "\n\n".join(s["text"] for s in page_segments)
    log.info(f"  ✓ {filename} — {len(full_text.split())} words, {total_pages} pages")

    return {
        "markdown":      full_text,
        "page_segments": page_segments,
        "tables":        tables,
        "metadata": {
            "source":       filename,
            "filepath":     str(p.resolve()),
            "doc_type":     doc_type,
            "headings":     [],
            "table_count":  0,
            "extracted_at": datetime.now().isoformat(),
        },
    }


def table_to_markdown(table) -> str:
    rows_data = []
    for row in table.rows:
        row_cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
        rows_data.append(row_cells)
    if not rows_data:
        return ""
    
    # Clean up empty rows
    rows_data = [r for r in rows_data if any(val for val in r)]
    if not rows_data:
        return ""

    # Build markdown table
    col_widths = [max(len(cell) for cell in col) for col in zip(*rows_data)]
    markdown_lines = []
    
    # Header
    header = "| " + " | ".join(val.ljust(w) for val, w in zip(rows_data[0], col_widths)) + " |"
    markdown_lines.append(header)
    
    # Separator
    separator = "|-" + "-|-".join("-" * w for w in col_widths) + "-|"
    markdown_lines.append(separator)
    
    # Body
    for row in rows_data[1:]:
        body_row = "| " + " | ".join(val.ljust(w) for val, w in zip(row, col_widths)) + " |"
        markdown_lines.append(body_row)
        
    return "\n".join(markdown_lines)


def format_paragraph(paragraph) -> str:
    text = paragraph.text.strip()
    if not text:
        return ""
    try:
        style_name = paragraph.style.name.lower() if paragraph.style else ""
    except Exception:
        style_name = ""
    if "bullet" in style_name:
        return f"- {text}"
    elif "number" in style_name:
        return f"1. {text}"
    return text


def extract_docx(filepath: str) -> dict:
    from docx import Document as DocxDocument
    from docx.document import Document as _DocxDocType
    from docx.table import Table as _DocxTable, _Cell
    from docx.text.paragraph import Paragraph as _DocxParagraph
    from docx.oxml.text.paragraph import CT_P
    from docx.oxml.table import CT_Tbl

    p = Path(filepath)
    filename = p.name
    doc_type = detect_doc_type(filename)
    log.info(f"Extracting DOCX: {filename}")

    doc = DocxDocument(str(p))
    
    def iter_block_items(parent):
        if isinstance(parent, _DocxDocType):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            raise TypeError("Unsupported parent type")
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield _DocxParagraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield _DocxTable(child, parent)

    current_section = ""
    page_segments = []
    tables = []
    current_segment_text = []
    
    for item in iter_block_items(doc):
        if isinstance(item, _DocxParagraph):
            text = format_paragraph(item)
            if not text:
                continue
            if item.style and item.style.name.startswith("Heading"):
                if current_segment_text:
                    page_segments.append({
                        "text": "\n".join(current_segment_text),
                        "page": 1,
                        "section": current_section
                    })
                    current_segment_text = []
                current_section = item.text.strip()
            current_segment_text.append(text)
        elif isinstance(item, _DocxTable):
            table_md = table_to_markdown(item)
            if table_md:
                if current_segment_text:
                    page_segments.append({
                        "text": "\n".join(current_segment_text),
                        "page": 1,
                        "section": current_section
                    })
                    current_segment_text = []
                tables.append({
                    "text": table_md,
                    "page": 1,
                    "table_index": len(tables),
                    "section": current_section
                })
                
    if current_segment_text:
        page_segments.append({
            "text": "\n".join(current_segment_text),
            "page": 1,
            "section": current_section
        })
        
    full_text = "\n\n".join(s["text"] for s in page_segments)
    
    return {
        "markdown": full_text,
        "page_segments": page_segments,
        "tables": tables,
        "metadata": {
            "source": filename,
            "filepath": str(p.resolve()),
            "doc_type": doc_type,
            "headings": [],
            "table_count": len(tables),
            "extracted_at": datetime.now().isoformat(),
        }
    }


def extract_pptx(filepath: str) -> dict:
    from pptx import Presentation
    p = Path(filepath)
    filename = p.name
    doc_type = detect_doc_type(filename)
    log.info(f"Extracting PPTX: {filename}")
    
    prs = Presentation(str(p))
    page_segments = []
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        slide_title = f"Slide {slide_idx}"
        
        title_shape = None
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_text = title_shape.text.strip()
            if title_text:
                slide_title = title_text
                slide_texts.append(title_text)
                
        for shape in slide.shapes:
            if shape == title_shape:
                continue
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
                
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            notes_text = notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(f"Speaker Notes: {notes_text}")
                
        if slide_texts:
            slide_content = "\n".join(slide_texts)
            page_segments.append({
                "text": slide_content,
                "page": slide_idx,
                "slide": slide_idx,
                "section": slide_title
            })
            
    full_text = "\n\n".join(s["text"] for s in page_segments)
    
    return {
        "markdown": full_text,
        "page_segments": page_segments,
        "tables": [],
        "metadata": {
            "source": filename,
            "filepath": str(p.resolve()),
            "doc_type": doc_type,
            "headings": [],
            "table_count": 0,
            "extracted_at": datetime.now().isoformat(),
        }
    }


def extract_xlsx(filepath: str) -> dict:
    import openpyxl
    p = Path(filepath)
    filename = p.name
    doc_type = detect_doc_type(filename)
    log.info(f"Extracting XLSX: {filename}")
    
    wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
    page_segments = []
    tables = []
    
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        
        markdown_rows = []
        consecutive_empty = 0
        max_cols = 0
        
        for r in sheet.iter_rows(values_only=True):
            if all(val is None or str(val).strip() == "" for val in r):
                consecutive_empty += 1
                if consecutive_empty > 20:
                    break
                continue
            consecutive_empty = 0
            
            # Find the last non-empty column index in this row to prune empty columns
            row_len = len(r)
            row_max_col = 0
            for idx in range(row_len - 1, -1, -1):
                if r[idx] is not None and str(r[idx]).strip() != "":
                    row_max_col = idx + 1
                    break
            max_cols = max(max_cols, row_max_col)
            
            row_vals = ["" if val is None else str(val).replace("\n", " ").strip() for val in r]
            markdown_rows.append(row_vals)
            
        if not markdown_rows:
            continue
            
        cropped_rows = [row[:max_cols] for row in markdown_rows]
        if not cropped_rows or not cropped_rows[0]:
            continue
            
        col_widths = [max(len(cell) for cell in col) for col in zip(*cropped_rows)]
        if not col_widths:
            continue
            
        markdown_lines = []
        header = "| " + " | ".join(val.ljust(w) for val, w in zip(cropped_rows[0], col_widths)) + " |"
        markdown_lines.append(header)
        separator = "|-" + "-|-".join("-" * w for w in col_widths) + "-|"
        markdown_lines.append(separator)
        for row in cropped_rows[1:]:
            body_row = "| " + " | ".join(val.ljust(w) for val, w in zip(row, col_widths)) + " |"
            markdown_lines.append(body_row)
            
        sheet_text = f"Worksheet: {sheet_name}\n\n" + "\n".join(markdown_lines)
        page_segments.append({
            "text": sheet_text,
            "page": 1,
            "worksheet": sheet_name,
            "section": f"Sheet: {sheet_name}"
        })
        
        tables.append({
            "text": "\n".join(markdown_lines),
            "page": 1,
            "table_index": len(tables),
            "section": f"Sheet: {sheet_name}",
            "worksheet": sheet_name
        })
        
    full_text = "\n\n".join(s["text"] for s in page_segments)
    wb.close()
    
    return {
        "markdown": full_text,
        "page_segments": page_segments,
        "tables": tables,
        "metadata": {
            "source": filename,
            "filepath": str(p.resolve()),
            "doc_type": doc_type,
            "headings": [],
            "table_count": len(tables),
            "extracted_at": datetime.now().isoformat(),
        }
    }


def extract_txt(filepath: str) -> dict:
    p = Path(filepath)
    filename = p.name
    doc_type = detect_doc_type(filename)
    log.info(f"Extracting TXT: {filename}")
    
    encodings = ["utf-8", "latin-1", "cp1252", "utf-16"]
    text = ""
    for enc in encodings:
        try:
            text = p.read_text(encoding=enc)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = p.read_bytes().decode("utf-8", errors="replace")
        
    text = clean_text(text)
    
    page_segments = [{
        "text": text,
        "page": 1,
        "section": "Plain Text Content"
    }]
    
    return {
        "markdown": text,
        "page_segments": page_segments,
        "tables": [],
        "metadata": {
            "source": filename,
            "filepath": str(p.resolve()),
            "doc_type": doc_type,
            "headings": [],
            "table_count": 0,
            "extracted_at": datetime.now().isoformat(),
        }
    }

def extract_zip(filepath: str) -> dict:
    """Extract ZIP archive and process supported files inside.
    Returns a combined extraction dict merging results of all extracted files.
    """
    p = Path(filepath)
    filename = p.name
    log.info(f"Extracting ZIP: {filename}")
    # Create temporary directory for extraction
    with tempfile.TemporaryDirectory() as tmpdir:
        with zipfile.ZipFile(str(p), 'r') as z:
            z.extractall(tmpdir)
        # Gather all supported files
        extracted_markdown = []
        page_segments = []
        tables = []
        total_tables = 0
        for inner_path in Path(tmpdir).rglob('*'):
            if inner_path.is_file():
                ext = inner_path.suffix.lower()
                if ext in SUPPORTED_EXTENSIONS and ext != '.zip':
                    try:
                        result = extract_file(str(inner_path))
                        extracted_markdown.append(result.get('markdown', ''))
                        page_segments.extend(result.get('page_segments', []))
                        tables.extend(result.get('tables', []))
                        total_tables += len(result.get('tables', []))
                    except Exception as e:
                        log.warning(f"Failed to extract {inner_path}: {e}")
        combined_md = "\n\n".join(extracted_markdown)
        return {
            "markdown": combined_md,
            "page_segments": page_segments,
            "tables": tables,
            "metadata": {
                "source": filename,
                "filepath": str(p.resolve()),
                "doc_type": "zip",
                "inner_files": [str(p.relative_to(tmpdir)) for p in Path(tmpdir).rglob('*') if p.is_file()],
                "table_count": total_tables,
                "extracted_at": datetime.now().isoformat(),
            },
        }



# ══════════════════════════════════════════════════════════════════════════════
#  MAIN EXTRACTION ROUTER
# ══════════════════════════════════════════════════════════════════════════════

def extract_file(filepath: str) -> dict:
    """
    Route to the correct extractor based on file extension.
    All extractors return the same dict shape:
      {markdown, page_segments, tables, metadata}
    """
    ext = Path(filepath).suffix.lower()

    if ext == ".pdf":
        return extract_pdf(filepath)
    elif ext == ".docx":
        return extract_docx(filepath)
    elif ext == ".pptx":
        return extract_pptx(filepath)
    elif ext == ".xlsx":
        return extract_xlsx(filepath)
    elif ext == ".txt":
        return extract_txt(filepath)
    elif ext == ".zip":
        return extract_zip(filepath)
# duplicate zip case removed
    else:
        raise ValueError(
            f"Unsupported file type: {ext}. "
            f"Supported: {sorted(SUPPORTED_EXTENSIONS)}"
        )


# ══════════════════════════════════════════════════════════════════════════════
#  CHUNKING — page-aware
# ══════════════════════════════════════════════════════════════════════════════

def chunk_extraction(extracted: dict) -> list[Document]:
    """
    Chunk extracted content into LangChain Documents.
    Preserves real page numbers and section headings per chunk.
    """
    meta          = extracted["metadata"]
    cfg           = CHUNK_CONFIG.get(meta["doc_type"], CHUNK_CONFIG["unknown"])
    docs: list[Document] = []
    seen: set[str]       = set()

    char_splitter = RecursiveCharacterTextSplitter(
        chunk_size=cfg["chunk_size"],
        chunk_overlap=cfg["chunk_overlap"],
        separators=["\n\n", "\n", ". ", "! ", "? ", " ", ""],
    )

    page_segments = extracted.get("page_segments", [])

    if page_segments:
        # ── Page-aware path (PDF via PyMuPDF) ────────────────────────────────
        for seg in page_segments:
            seg_text    = seg.get("text", "").strip()
            seg_page    = seg.get("page", 1)
            seg_section = seg.get("section", "")

            if len(seg_text.split()) < 5:
                continue

            sub_chunks = (
                char_splitter.split_text(seg_text)
                if len(seg_text) > cfg["chunk_size"]
                else [seg_text]
            )

            for text in sub_chunks:
                text = text.strip()
                if len(text.split()) < 8:
                    continue
                # FIX BUG 3: skip Table-of-Contents chunks — they match queries
                # by topic name but contain zero instructional content.
                if _is_toc_chunk(text):
                    continue
                h = hashlib.md5(text.encode()).hexdigest()
                if h in seen:
                    continue
                seen.add(h)

                idx = len(docs)
                docs.append(Document(
                    page_content=text,
                    metadata={
                        "chunk_id":     chunk_id(text, meta["source"], idx),
                        "source":       meta["source"],
                        "filepath":     meta["filepath"],
                        "doc_type":     meta["doc_type"],
                        "section":      seg_section,
                        "page":         seg_page,
                        "slide":        seg.get("slide"),
                        "worksheet":    seg.get("worksheet"),
                        "language":     detect_language(text),
                        "chunk_index":  idx,
                        "word_count":   len(text.split()),
                        "is_table":     False,
                        "processed_at": meta["extracted_at"],
                    }
                ))
    else:
        # ── Fallback: markdown header splitting ───────────────────────────────
        try:
            header_splitter = MarkdownHeaderTextSplitter(
                headers_to_split_on=_MD_HEADERS, strip_headers=False
            )
            header_chunks = header_splitter.split_text(extracted["markdown"])
        except Exception:
            header_chunks = [Document(page_content=extracted["markdown"], metadata={})]

        for h_chunk in header_chunks:
            for text in char_splitter.split_text(h_chunk.page_content):
                text = text.strip()
                if len(text.split()) < 8:
                    continue
                h = hashlib.md5(text.encode()).hexdigest()
                if h in seen:
                    continue
                seen.add(h)

                idx     = len(docs)
                section = (
                    h_chunk.metadata.get("h1")
                    or h_chunk.metadata.get("h2")
                    or h_chunk.metadata.get("h3", "")
                )
                docs.append(Document(
                    page_content=text,
                    metadata={
                        "chunk_id":     chunk_id(text, meta["source"], idx),
                        "source":       meta["source"],
                        "filepath":     meta["filepath"],
                        "doc_type":     meta["doc_type"],
                        "section":      section,
                        "page":         idx + 1,
                        "slide":        None,
                        "worksheet":    None,
                        "language":     detect_language(text),
                        "chunk_index":  idx,
                        "word_count":   len(text.split()),
                        "is_table":     False,
                        "processed_at": meta["extracted_at"],
                    }
                ))

    # ── Tables ────────────────────────────────────────────────────────────────
    for tbl in extracted.get("tables", []):
        text = tbl["text"].strip()
        if not text or len(text.split()) < 4:
            continue
        h = hashlib.md5(text.encode()).hexdigest()
        if h in seen:
            continue
        seen.add(h)
        idx = len(docs)
        docs.append(Document(
            page_content=text,
            metadata={
                "chunk_id":     chunk_id(text, meta["source"], idx),
                "source":       meta["source"],
                "filepath":     meta["filepath"],
                "doc_type":     meta["doc_type"],
                "section":      f"Table {tbl['table_index'] + 1}",
                "page":         tbl.get("page", 1),
                "slide":        tbl.get("slide"),
                "worksheet":    tbl.get("worksheet"),
                "language":     detect_language(text),
                "chunk_index":  idx,
                "word_count":   len(text.split()),
                "is_table":     True,
                "processed_at": meta["extracted_at"],
            }
        ))

    log.info(f"  Chunked {meta['source']} → {len(docs)} chunks")
    return docs


def chunk_all_extractions(extractions: list[dict]) -> list[Document]:
    all_docs: list[Document] = []
    for ex in extractions:
        all_docs.extend(chunk_extraction(ex))
    log.info(f"Total chunks: {len(all_docs)}")
    return all_docs


# ══════════════════════════════════════════════════════════════════════════════
#  SAVE — handoff to member2
# ══════════════════════════════════════════════════════════════════════════════

def save_chunks_json(documents: list[Document], path: str = CHUNKS_JSON_PATH) -> None:
    Path(path).parent.mkdir(parents=True, exist_ok=True)

    # Load existing chunks
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            payload = json.load(f)
    else:
        payload = []

    # Add only new chunks
    new_chunks = [
        {
            "text": d.page_content,
            "page": d.metadata.get("page", d.metadata.get("chunk_index", 0)),
            "slide": d.metadata.get("slide"),
            "worksheet": d.metadata.get("worksheet"),
            "source_file": d.metadata.get("source", ""),
            "section": d.metadata.get("section", ""),
            "doc_type": d.metadata.get("doc_type", "unknown"),
            "is_table": d.metadata.get("is_table", False),
            "chunk_index": d.metadata.get("chunk_index", 0),
        }
        for d in documents
    ]

    payload.extend(new_chunks)

    with open(path, "w", encoding="utf-8") as f:
        json.dump(payload, f, indent=2, ensure_ascii=False)

    log.info(f"Added {len(new_chunks)} new chunks. Total chunks: {len(payload)}")

# ══════════════════════════════════════════════════════════════════════════════
#  PUBLIC API
# ══════════════════════════════════════════════════════════════════════════════

def process_documents(file_paths: list[str]) -> list[Document]:
    """
    Full pipeline: extract → chunk → save.
    Called by api.py background task.
    Returns list[Document] ready for member2 FAISS indexing.
    """
    extractions = []

    for fp in file_paths:
        try:
            extractions.append(extract_file(fp))
        except Exception as e:
            log.error(f"Failed to extract {fp}: {e}")

    if not extractions:
        return []

    documents = chunk_all_extractions(extractions)
    save_chunks_json(documents)

    return documents


# ══════════════════════════════════════════════════════════════════════════════
#  HTML PREVIEW GENERATORS FOR IN-BROWSER VIEWING
# ══════════════════════════════════════════════════════════════════════════════

def docx_to_html(filepath: Path) -> str:
    from docx import Document as DocxDocument
    from docx.document import Document as _DocxDocType
    from docx.table import Table as _DocxTable, _Cell
    from docx.text.paragraph import Paragraph as _DocxParagraph
    from docx.oxml.text.paragraph import CT_P
    from docx.oxml.table import CT_Tbl

    doc = DocxDocument(str(filepath))

    def iter_block_items(parent):
        if isinstance(parent, _DocxDocType):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            raise TypeError("Unsupported parent type")
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield _DocxParagraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield _DocxTable(child, parent)

    html_parts = []
    
    html_parts.append("""
    <style>
        .preview-body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Helvetica, Arial, sans-serif; color: #1f2937; line-height: 1.6; max-width: 900px; margin: 0 auto; }
        .preview-body h1, .preview-body h2, .preview-body h3, .preview-body h4 { color: #111827; font-weight: 600; margin-top: 1.5em; margin-bottom: 0.5em; }
        .preview-body h1 { border-bottom: 1px solid #e5e7eb; padding-bottom: 0.3em; font-size: 1.8em; }
        .preview-body h2 { font-size: 1.4em; }
        .preview-body h3 { font-size: 1.2em; }
        .preview-body p { margin-top: 0; margin-bottom: 1em; }
        .preview-body ul, .preview-body ol { margin-top: 0; margin-bottom: 1em; padding-left: 2em; }
        .preview-body li { margin-bottom: 0.25em; }
        .preview-body table { border-collapse: collapse; width: 100%; margin: 1.5em 0; font-size: 0.9em; }
        .preview-body th, .preview-body td { border: 1px solid #e5e7eb; padding: 8px 12px; text-align: left; }
        .preview-body th { background-color: #f9fafb; font-weight: 600; }
        .preview-body .highlight-section { background-color: #f3e8ff; border-left: 4px solid #9333ea; padding: 8px; margin: 10px 0; border-radius: 0 4px 4px 0; animation: pulse-highlight 2s ease-in-out; }
        @keyframes pulse-highlight {
            0% { background-color: #f3e8ff; }
            50% { background-color: #e9d5ff; }
            100% { background-color: #f3e8ff; }
        }
    </style>
    <div class="preview-body">
    """)

    def slugify(text: str) -> str:
        import re
        return "heading-" + re.sub(r'[^a-z0-9]+', '-', text.lower()).strip('-')

    in_list = False
    list_type = None

    for item in iter_block_items(doc):
        if isinstance(item, _DocxParagraph):
            text = item.text.strip()
            if not text:
                continue
            
            style_name = item.style.name.lower() if item.style else ""
            is_bullet = "bullet" in style_name
            is_number = "number" in style_name
            
            if is_bullet or is_number:
                current_type = "ul" if is_bullet else "ol"
                if not in_list:
                    html_parts.append(f"<{current_type}>")
                    in_list = True
                    list_type = current_type
                elif list_type != current_type:
                    html_parts.append(f"</{list_type}>")
                    html_parts.append(f"<{current_type}>")
                    list_type = current_type
                html_parts.append(f"<li>{text}</li>")
            else:
                if in_list:
                    html_parts.append(f"</{list_type}>")
                    in_list = False
                    list_type = None
                
                if item.style and style_name.startswith("heading"):
                    level = 1
                    try:
                        level_match = re.search(r'\d', style_name)
                        if level_match:
                            level = int(level_match.group())
                    except:
                        pass
                    h_tag = f"h{min(max(level, 1), 6)}"
                    heading_id = slugify(text)
                    html_parts.append(f"<{h_tag} id=\"{heading_id}\">{text}</{h_tag}>")
                else:
                    html_parts.append(f"<p>{text}</p>")
                    
        elif isinstance(item, _DocxTable):
            if in_list:
                html_parts.append(f"</{list_type}>")
                in_list = False
                list_type = None
                
            html_parts.append("<table>")
            for r_idx, row in enumerate(item.rows):
                html_parts.append("<tr>")
                for cell in row.cells:
                    cell_tag = "th" if r_idx == 0 else "td"
                    html_parts.append(f"<{cell_tag}>{cell.text.strip()}</{cell_tag}>")
                html_parts.append("</tr>")
            html_parts.append("</table>")
            
    if in_list:
        html_parts.append(f"</{list_type}>")

    html_parts.append("</div>")
    return "\n".join(html_parts)


def pptx_to_html(filepath: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(filepath))
    
    html_parts = []
    html_parts.append("""
    <style>
        .preview-body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Helvetica, Arial, sans-serif; color: #1f2937; line-height: 1.6; max-width: 900px; margin: 0 auto; padding-bottom: 50px; }
        .slide { background-color: #ffffff; border: 1px solid #e5e7eb; border-radius: 8px; box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1); margin-bottom: 30px; padding: 40px; aspect-ratio: 16/9; position: relative; display: flex; flex-direction: column; justify-content: flex-start; box-sizing: border-box; }
        .slide-header { border-bottom: 2px solid #9333ea; color: #111827; font-size: 1.6em; font-weight: 600; margin-bottom: 20px; padding-bottom: 10px; }
        .slide-content { flex: 1; font-size: 1.1em; color: #374151; overflow-y: auto; }
        .slide-notes { background-color: #fdf4ff; border-left: 4px solid #d8b4fe; border-radius: 0 4px 4px 0; font-size: 0.9em; font-style: italic; margin-top: 15px; padding: 8px 12px; color: #581c87; }
        .slide-number { position: absolute; bottom: 15px; right: 20px; font-size: 0.8em; color: #9ca3af; font-weight: 500; }
        .highlight-section { border: 3px solid #a855f7; box-shadow: 0 0 15px rgba(168, 85, 247, 0.4); animation: slide-border-pulse 2s infinite; }
        @keyframes slide-border-pulse {
            0% { box-shadow: 0 0 15px rgba(168, 85, 247, 0.4); }
            50% { box-shadow: 0 0 25px rgba(168, 85, 247, 0.7); }
            100% { box-shadow: 0 0 15px rgba(168, 85, 247, 0.4); }
        }
    </style>
    <div class="preview-body">
    """)
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_title = f"Slide {slide_idx}"
        slide_texts = []
        
        title_shape = None
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_text = title_shape.text.strip()
            if title_text:
                slide_title = title_text
                
        for shape in slide.shapes:
            if shape == title_shape:
                continue
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
                
        notes_text = ""
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            notes_text = notes_slide.notes_text_frame.text.strip()
            
        html_parts.append(f'<div class="slide" id="slide-{slide_idx}">')
        html_parts.append(f'  <div class="slide-header">{slide_title}</div>')
        html_parts.append('  <div class="slide-content">')
        for txt in slide_texts:
            html_parts.append(f'    <p>{txt}</p>')
        html_parts.append('  </div>')
        if notes_text:
            html_parts.append(f'  <div class="slide-notes"><strong>Speaker Notes:</strong> {notes_text}</div>')
        html_parts.append(f'  <div class="slide-number">Slide {slide_idx} of {len(prs.slides)}</div>')
        html_parts.append('</div>')
        
    html_parts.append("</div>")
    return "\n".join(html_parts)


def xlsx_to_html(filepath: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(filepath), data_only=True, read_only=True)
    
    html_parts = []
    html_parts.append("""
    <style>
        .preview-body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Helvetica, Arial, sans-serif; color: #1f2937; padding: 10px; }
        .tabs { display: flex; border-bottom: 2px solid #e5e7eb; margin-bottom: 20px; flex-wrap: wrap; }
        .tab-btn { padding: 8px 16px; font-weight: 500; cursor: pointer; border: none; background: none; border-bottom: 2px solid transparent; color: #6b7280; font-size: 0.9em; }
        .tab-btn.active { border-bottom-color: #9333ea; color: #9333ea; font-weight: 600; }
        .sheet-container { display: none; overflow-x: auto; }
        .sheet-container.active { display: block; }
        table { border-collapse: collapse; width: 100%; margin: 15px 0; font-size: 0.85em; }
        th, td { border: 1px solid #e5e7eb; padding: 6px 10px; text-align: left; min-width: 80px; }
        th { background-color: #f9fafb; font-weight: 600; }
    </style>
    <div class="preview-body">
    """)
    
    html_parts.append('<div class="tabs">')
    for sheet_idx, sheet_name in enumerate(wb.sheetnames):
        active_class = "active" if sheet_idx == 0 else ""
        html_parts.append(f'<button class="tab-btn {active_class}" onclick="showSheet(\'{sheet_name}\', event)">{sheet_name}</button>')
    html_parts.append('</div>')
    
    for sheet_idx, sheet_name in enumerate(wb.sheetnames):
        active_class = "active" if sheet_idx == 0 else ""
        html_parts.append(f'<div class="sheet-container {active_class}" id="sheet-{sheet_name}">')
        html_parts.append('  <table>')
        
        sheet = wb[sheet_name]
        consecutive_empty = 0
        for r_idx, row in enumerate(sheet.iter_rows(values_only=True)):
            if all(val is None or str(val).strip() == "" for val in row):
                consecutive_empty += 1
                if consecutive_empty > 20:
                    break
                continue
            consecutive_empty = 0
            
            html_parts.append('    <tr>')
            for val in row:
                cell_tag = "th" if r_idx == 0 else "td"
                val_str = "" if val is None else str(val).strip()
                html_parts.append(f'      <{cell_tag}>{val_str}</{cell_tag}>')
            html_parts.append('    </tr>')
            
        html_parts.append('  </table>')
        html_parts.append('</div>')
        
    html_parts.append("""
    </div>
    <script>
        function showSheet(name, event) {
            document.querySelectorAll('.sheet-container').forEach(el => el.classList.remove('active'));
            document.querySelectorAll('.tab-btn').forEach(el => el.classList.remove('active'));
            document.getElementById('sheet-' + name).classList.add('active');
            event.target.classList.add('active');
        }
    </script>
    """)
    wb.close()
    return "\n".join(html_parts)


# ══════════════════════════════════════════════════════════════════════════════
#  CLI
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import argparse
    logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")

    parser = argparse.ArgumentParser(description="Member 1 — Extract & Chunk Documents")
    parser.add_argument("--input", "-i", required=True, help="PDF file or directory")
    args = parser.parse_args()

    docs = process_documents(
        [str(f) for f in Path(args.input).rglob("*") if f.suffix.lower() in SUPPORTED_EXTENSIONS]
        if Path(args.input).is_dir()
        else [args.input]
    )
    print(f"\n✅ Done — {len(docs)} chunks ready in {CHUNKS_JSON_PATH}")